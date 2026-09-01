"""Assemble the Week 2 Tuesday deck from scripts/demo_photos/.

Run every gen_*.py first, then:  python build_deck.py
Output: scripts/W02_How_Text_Becomes_Numbers.pptx

Slides are 16:9, dark (#1b1e26) to match the GIF backgrounds and the course
site's slate theme. GIFs animate in slideshow mode. Each slide carries speaker
notes naming the study-guide concepts it teaches, so the deck and
docs/study-guide.md stay in step.
"""

import glob
import os
import sys

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
sys.path.insert(0, SCRIPT_DIR)

from deck_order import ORDER  # noqa: E402

PHOTOS = os.path.join(SCRIPT_DIR, "demo_photos")
OUT = os.path.join(SCRIPT_DIR, "W02_How_Text_Becomes_Numbers.pptx")
BG = RGBColor(0x1B, 0x1E, 0x26)


# Speaker notes, keyed by slide file prefix. First line = the beat to hit;
# "SG:" names the study-guide entries; "Ask:" is a question to throw out.
NOTES = {
    "w02_s01": (
        "Cold open: this map is real — every dot is a word placed by arithmetic. "
        "By the end of today they can read this picture.\n"
        "Session: How Text Becomes Numbers (w02-tue)."
    ),
    "w02_s02": (
        "The entire lecture is the ? box. Models never see letters — only numbers.\n"
        "SG: vector · embedding.\n"
        "Ask: what do YOU think is in the box?"
    ),
    "w02_s03": (
        "Tokens are not words. The primer's own example — 'coders' becomes "
        "cod+ers. Strawberry lands here: it's not that models can't count, "
        "they never saw the letters.\n"
        "SG: token-and-tokenization · vocabulary.\n"
        "Thursday's lab: we call the real token counter on their own text."
    ),
    "w02_s04": (
        "First honest attempt. Let them compute cat·dog aloud — it's zero, and "
        "zero for EVERY pair. Identity without meaning, 50k slots of zeros.\n"
        "SG: one-hot-encoding · sparse-and-dense-vectors · dot-product (preview)."
    ),
    "w02_s04b": (
        "DISCUSSION (3-5 min). Fish for: (1) similarity should be graded, "
        "not binary — they'll invent cosine before meeting it; (2) nobody "
        "hand-decides — it has to come from data, cheaply, at scale; "
        "(3) far fewer slots, densely used — they'll invent dimensionality.\n"
        "Don't resolve the questions. The field didn't either: next comes "
        "what everyone did first — count. Their answers get delivered at "
        "the scorecard, after counting hits its ceiling."
    ),
    "w02_s05": (
        "The bridge intuition (from Cornell's CS4782 deck): an embedding is a "
        "scorecard. Similar words = similar rows. Then the twist: nobody writes "
        "it — training fills it in, axes unlabelled, ~1,000 columns.\n"
        "SG: embedding · dimensionality.\n"
        "Ask: what would a 'verb-ness' column score for 'run'?"
    ),
    "w02_s05b": (
        "The scorecard, earned. Derive ONE column live: count how often each "
        "word appears near feline company (purr, fur, claws), rescale, and a "
        "column materialises. king scores low — no cat-words in his company.\n"
        "SG: embedding · distributional-hypothesis · corpus.\n"
        "Key line: training runs this for hundreds of columns at once and "
        "never names any of them — which is why axes are unlabelled."
    ),
    "w02_s06": (
        "Watch meaning become geometry. Neighbour list on the right is how anyone "
        "actually inspects a space.\n"
        "SG: embedding-space · nearest-neighbours.\n"
        "Caveat now, not later: good/bad are neighbours — similar ≠ synonym."
    ),
    "w02_s06b": (
        "Flatland escape. The opening frame IS the 2-D plot they've been "
        "looking at — top-down. Tilt: a third axis was always there, and king/"
        "coffee, near-neighbours in 2-D, fly apart. Let the rotation run.\n"
        "Then the jump: GPT-2 768 dims, DeepSeek-R1 7,168.\n"
        "SG: dimensionality · embedding-space · dimensionality-reduction.\n"
        "Pros/cons: room for distinctions vs memory+compute, overfitting, "
        "and unviewability — every plot from here on is a projection.\n"
        "Ask: why not just use 10 dimensions? why not a million?"
    ),
    "w02_s07": (
        "The angle IS the similarity. Watch the meter as the arrow swings. "
        "Length changes nothing — that's the whole point of cosine over dot.\n"
        "SG: dot-product · cosine-similarity.\n"
        "This number runs every retrieval system they'll build in Week 10."
    ),
    "w02_s08": (
        "The famous demo, then its cracks. Run it live in the projector "
        "(projector.tensorflow.org) right after this slide.\n"
        "SG: vector-arithmetic-and-its-limits.\n"
        "In-class activity: find one analogy that works and one that fails."
    ),
    "w02_s09": (
        "THE PIVOT. Counting just hit its wall; this is the question that "
        "gets past it. 'zorp': they infer a drink from three sentences — they just "
        "ran the distributional hypothesis themselves. Firth quote lands here.\n"
        "SG: distributional-hypothesis · corpus.\n"
        "This is the hinge of the whole week — everything learned rests on it."
    ),
    "w02_s10": (
        "Rung 1. Order dies in the bag: dog-bites-man == man-bites-dog, and no "
        "setting fixes it — the method is working as designed.\n"
        "SG: bag-of-words · count-vectorization.\n"
        "Still the backbone of keyword search everywhere."
    ),
    "w02_s11": (
        "Rung 2. TF x IDF; 'the' is deleted by arithmetic, not by a list. "
        "Fully interpretable — every score explainable. Worth admiring.\n"
        "SG: term-frequency · inverse-document-frequency · tf-idf · stop-words."
    ),
    "w02_s12": (
        "Fast tour, 90 seconds. The trap to name: stop-word removal deletes "
        "'not' — negation lives in stop words.\n"
        "SG: stemming · lemmatization · stop-words · n-grams.\n"
        "Each knob is an ADR-sized decision, not a default."
    ),
    "w02_s12b": (
        "DISCUSSION (3-5 min). Fish for: (1) no — counters can't see "
        "similarity, only co-occurrence in documents; (2) nothing returned — "
        "the synonym gap, this is why search moved to embeddings; (3) the "
        "bag is identical, bigrams patch it locally and explode the "
        "vocabulary.\n"
        "Next slide is the blob: counting's ceiling in real data. Then "
        "zorp — the question that escapes it — so question 1 gets its "
        "answer two slides from now. Don't give it away here."
    ),
    "w02_s13": (
        "The ceiling. Real TF-IDF spaces collapse into a blob (the HF primer "
        "shows this with PCA). Counting cannot learn cat≈kitten. A different "
        "QUESTION was needed — that pivot is the next slide.\n"
        "SG: tf-idf (pitfall) · dimensionality-reduction."
    ),
    "w02_s14": (
        "The new question: predict your neighbours. Window slides, pairs fall "
        "out — billions of free training examples, no human labels.\n"
        "SG: sliding-window · skip-gram · cbow.\n"
        "Skip-gram: word→context. CBOW: context→word. Both same bet."
    ),
    "w02_s15": (
        "word2vec's secret: the network is a pretext. After training, throw the "
        "output layer away — matrix W (one row per word) was the treasure.\n"
        "SG: word2vec · pretext-task · neural-network-and-layers · "
        "parameters-and-weights · softmax.\n"
        "Same trick later: BERT's masking, GPT's next-token."
    ),
    "w02_s16": (
        "Full softmax over 50k words per step = untrainable. Negative sampling "
        "changes the question to yes/no with a few random fakes.\n"
        "SG: negative-sampling · softmax · training-and-inference.\n"
        "Name the pattern: cheap approximation beats exact computation."
    ),
    "w02_s16b": (
        "DISCUSSION (3-5 min). Fish for: (1) an average of every sense — "
        "a point between rivers and money, wrong about both; (2) long-range "
        "reference (the trophy/suitcase 'it') — windows can't reach it; "
        "(3) averaging word vectors = a bag of vectors — they just "
        "reinvented bag-of-words one level up. Enjoy that landing.\n"
        "Next slide is the bank crime scene they just predicted."
    ),
    "w02_s17": (
        "The flaw that ends the word2vec era: 'bank' gets ONE vector, wrong "
        "about both senses. Static vs contextual is discussion question 3.\n"
        "SG: polysemy · static-vs-contextual-embeddings."
    ),
    "w02_s18": (
        "The fix, honestly told: an LLM's embedding layer is STILL a lookup "
        "table. Context is added by the layers above — watch the two 'bank's "
        "diverge on the way up. The mixing machinery is attention = Week 3.\n"
        "SG: static-vs-contextual-embeddings · bert · encoder-and-decoder · "
        "masked-language-modelling."
    ),
    "w02_s18b": (
        "DISCUSSION (3-5 min). Fish for: (1) no more precomputed "
        "dictionary — every representation needs a forward pass: compute, "
        "cost, latency (Week 11 territory); (2) which words to listen to "
        "is exactly attention — question 2 IS next week's paper title; "
        "(3) bias, staleness, corpus-bounds — the mirror slide lands in "
        "two slides and they'll have predicted it.\n"
        "The arc closes here: every method today existed to fix the one "
        "before it — say that out loud before the recap."
    ),
    "w02_s19": (
        "Two receipts that the space mirrors its corpus: Bolukbasi's she–he "
        "occupations, and 'broadcast' drifting across centuries.\n"
        "SG: corpus · distributional-hypothesis (pitfall).\n"
        "Moral: interrogate the corpus, not the model's 'understanding'."
    ),
    "w02_s20": (
        "The whole week in one picture. Each rung fixed the last one's flaw and "
        "exposed its own. Reading responses can ask for any rung's ✓ and ✗.\n"
        "SG: bag-of-words · tf-idf · word2vec · bert (the ladder)."
    ),
    "w02_s21": (
        "Cliffhanger. Seq2seq squeezes a sentence through one fixed vector; "
        "long sentences die. The 2017 fix is attention — the paper this course "
        "is named after. Read Alammar's seq2seq post before Tuesday.\n"
        "SG: sequence-to-sequence · hidden-state · fixed-length-bottleneck · "
        "recurrent-neural-network."
    ),
}


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    ordered = [os.path.join(PHOTOS, f) for f in ORDER]
    missing = [f for f in ordered if not os.path.exists(f)]
    if missing:
        raise SystemExit(
            "missing media (run the gen scripts first): "
            + ", ".join(os.path.basename(f) for f in missing)
        )
    extra = sorted(
        os.path.basename(f)
        for f in glob.glob(os.path.join(PHOTOS, "w02_s*.*"))
        if not f.endswith("_preview.png") and os.path.basename(f) not in ORDER
    )
    if extra:
        raise SystemExit(f"media not listed in ORDER (add them): {extra}")
    files = ordered

    for path in files:
        slide = prs.slides.add_slide(blank)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = BG
        slide.shapes.add_picture(path, 0, 0, width=prs.slide_width, height=prs.slide_height)
        key = "_".join(os.path.basename(path).split("_")[:2])
        note = NOTES.get(key)
        if note:
            slide.notes_slide.notes_text_frame.text = note

    prs.save(OUT)
    print(f"saved {OUT} ({len(files)} slides, {os.path.getsize(OUT) / 1e6:.1f} MB)")


if __name__ == "__main__":
    main()
